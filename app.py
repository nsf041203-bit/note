import asyncio
import os
import traceback
import uuid
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Literal

import aiofiles
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, UploadFile
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

# 导入数据库函数
from database import (
    init_db, create_course, create_outline, parse_outline_to_sections,
    get_course_by_name, get_latest_outline, save_session, get_session,
    get_course_sessions, get_course_by_id, list_courses_with_stats,
    rename_course, update_course_color, delete_course, delete_session,
    get_course_materials, update_session_notes, get_outline_by_id,
    rename_session,
)

load_dotenv()

# ===== 运行环境与数据目录 =====
# Vercel 等 serverless 平台文件系统只读，仅 /tmp 可写。
# 通过 VERCEL 环境变量自动判断；本地运行时一切照旧（仓库根目录）。
IS_SERVERLESS = bool(os.getenv("VERCEL") or os.getenv("AWS_LAMBDA_FUNCTION_NAME"))
DATA_DIR = Path(os.getenv("DATA_DIR") or ("/tmp/notesai" if IS_SERVERLESS else "."))
# 数据目录必须先于数据库初始化创建（serverless 上 DATA_DIR 在 /tmp，子目录默认不存在）
DATA_DIR.mkdir(parents=True, exist_ok=True)

# 让数据库也落到可写目录（database.py 读取 KNOWLEDGE_DB）
if "KNOWLEDGE_DB" not in os.environ:
    os.environ["KNOWLEDGE_DB"] = str(DATA_DIR / "knowledge.db")

# 初始化数据库
init_db()

app = FastAPI(title="课程笔记自动整理系统")

# 静态资源与页面用基于本文件的绝对路径，避免 serverless 下工作目录不确定导致 404
BASE_DIR = Path(__file__).resolve().parent
STATIC_DIR = BASE_DIR / "static"
app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")

UPLOAD_DIR = DATA_DIR / "uploads"
OUTPUT_DIR = DATA_DIR / "outputs"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 任务状态存储 {job_id: {"status": ..., "progress": ..., "files": [...]}}
jobs: dict[str, dict] = {}


def _set(job_id: str, status: str, progress: str):
    jobs[job_id]["status"] = status
    jobs[job_id]["progress"] = progress


def _extract_docx_markdown(path: str) -> str:
    """将 Word 文档转成 Markdown，再交给 AI 整理。"""
    try:
        return _extract_docx_markdown_with_python_docx(path)
    except Exception as exc:
        print(f"[Word提取] python-docx 解析失败，改用 XML 兜底: {exc!r}")
        return _extract_docx_markdown_from_xml(path)


def _extract_docx_markdown_with_python_docx(path: str) -> str:
    """优先使用 python-docx，把普通段落、标题、列表和表格转成 Markdown。"""
    from docx import Document

    doc = Document(path)
    parts: list[str] = []

    for p in doc.paragraphs:
        md = _docx_paragraph_to_markdown(p)
        if md:
            parts.append(md)

    for table in doc.tables:
        table_md = _docx_table_to_markdown(table)
        if table_md:
            parts.append(table_md)

    try:
        for section in doc.sections:
            for block in (section.header.paragraphs, section.footer.paragraphs):
                for p in block:
                    md = _docx_paragraph_to_markdown(p)
                    if md:
                        parts.append(md)
    except Exception as exc:
        # 某些 Word 文件的页眉页脚关系不完整，正文仍可继续使用。
        print(f"[Word提取] 页眉页脚读取失败，已跳过: {exc!r}")

    return "\n\n".join(parts).strip()


def _docx_paragraph_to_markdown(paragraph) -> str:
    """把 python-docx 段落尽量转换成 Markdown。"""
    text = paragraph.text.strip()
    if not text:
        return ""

    style_name = (paragraph.style.name if paragraph.style else "").lower()
    if "heading" in style_name or "标题" in style_name:
        level = 2
        for token in ("1", "2", "3", "4", "5", "6"):
            if token in style_name:
                level = min(int(token), 6)
                break
        return f"{'#' * level} {text}"

    if "list bullet" in style_name or "项目符号" in style_name:
        return f"- {text}"
    if "list number" in style_name or "编号" in style_name:
        return f"1. {text}"

    return text


def _docx_table_to_markdown(table) -> str:
    """把 Word 表格转成 Markdown 表格。"""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [_clean_table_cell(cell.text) for cell in row.cells]
        if any(cells):
            rows.append(cells)

    if not rows:
        return ""

    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    header = rows[0]
    separator = ["---"] * width
    body = rows[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _clean_table_cell(text: str) -> str:
    return " ".join(text.replace("|", "\\|").split())


def _extract_docx_markdown_from_xml(path: str) -> str:
    """直接读取 docx 压缩包内 XML，作为 python-docx 失败时的 Markdown 兜底。"""
    paragraphs: list[str] = []
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    with zipfile.ZipFile(path) as zf:
        xml_names = [
            name for name in zf.namelist()
            if name == "word/document.xml"
            or name.startswith("word/header")
            or name.startswith("word/footer")
        ]
        for name in xml_names:
            root = ET.fromstring(zf.read(name))
            for para in root.findall(".//w:p", ns):
                chunks: list[str] = []
                for node in para.iter():
                    if node.tag == f"{{{ns['w']}}}t" and node.text:
                        chunks.append(node.text)
                    elif node.tag == f"{{{ns['w']}}}tab":
                        chunks.append("\t")
                    elif node.tag == f"{{{ns['w']}}}br":
                        chunks.append("\n")
                text = "".join(chunks).strip()
                if text:
                    paragraphs.append(text)

    return "\n\n".join(paragraphs).strip()


async def _run_pipeline(job_id: str, subject: str, outline: str, saved_files: list[Path], color: str = "blue"):
    """后台运行完整处理流程"""
    job_dir = OUTPUT_DIR / job_id
    job_dir.mkdir(exist_ok=True)
    texts: list[str] = []

    try:
        # 1. 提取各类文件文字
        _set(job_id, "running", "📂 正在分析上传文件...")
        for fp in saved_files:
            suffix = fp.suffix.lower()

            if suffix in {".mp3", ".wav", ".m4a", ".ogg", ".flac"}:
                _set(job_id, "running", f"🎙️ 语音转文字: {fp.name}")
                from modules.audio import transcribe
                texts.append(await asyncio.to_thread(transcribe, str(fp)))

            elif suffix in {".jpg", ".jpeg", ".png", ".bmp", ".webp"}:
                _set(job_id, "running", f"🖼️ OCR识别: {fp.name}")
                try:
                    from modules.ocr import extract_text
                except ModuleNotFoundError as exc:
                    if exc.name == "paddleocr":
                        _set(
                            job_id,
                            "error",
                            "❌ 当前部署环境未安装 PaddleOCR，暂时不能从图片中提取文字。请改用本地完整环境，或先把图片文字转成 TXT/Word/PDF 后再上传。",
                        )
                        return
                    raise
                texts.append(await asyncio.to_thread(extract_text, str(fp)))

            elif suffix in {".txt", ".md", ".srt"}:
                texts.append(fp.read_text(encoding="utf-8", errors="ignore"))

            elif suffix in {".docx"}:
                _set(job_id, "running", f"📄 Word转Markdown: {fp.name}")
                texts.append(await asyncio.to_thread(_extract_docx_markdown, str(fp)))

            elif suffix in {".pptx"}:
                _set(job_id, "running", f"📊 提取PPT文字: {fp.name}")
                from pptx import Presentation
                prs = Presentation(str(fp))
                slide_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            slide_texts.append(shape.text_frame.text)
                texts.append("\n".join(slide_texts))

            elif suffix == ".pdf":
                _set(job_id, "running", f"📄 提取PDF文字: {fp.name}")
                from pdfminer.high_level import extract_text as pdf_extract
                texts.append(await asyncio.to_thread(pdf_extract, str(fp)))

        if not texts or not any(t.strip() for t in texts):
            _set(job_id, "error", "❌ 未能从文件中提取到任何文字")
            return

        raw_text = "\n\n---\n\n".join(t for t in texts if t.strip())

        # 2. AI 结构化
        _set(job_id, "running", "🤖 AI正在整理笔记...")
        from modules.structurer import structure
        notes_md, mindmap_md = await asyncio.to_thread(structure, raw_text, subject, outline)

        # 3. 保存输出文件
        _set(job_id, "running", "📝 生成输出文件...")
        safe_name = subject.replace("/", "_").replace("\\", "_") or "课程笔记"
        from modules.exporter import to_markdown, to_docx
        # 注意：PDF 导出需要安装 GTK，见 README.md

        md_path = job_dir / f"{safe_name}.md"
        docx_path = job_dir / f"{safe_name}.docx"
        mindmap_path = job_dir / "mindmap.md"

        await asyncio.to_thread(to_markdown, notes_md, str(md_path))
        await asyncio.to_thread(to_docx, notes_md, str(docx_path), subject)
        mindmap_path.write_text(mindmap_md, encoding="utf-8")

        jobs[job_id]["files"] = [
            {"name": f"{safe_name}.md", "label": "Markdown 笔记"},
            {"name": f"{safe_name}.docx", "label": "Word 文档"},
        ]

        # 尝试生成 PDF（如果 GTK 已安装）
        try:
            from modules.exporter import to_pdf
            pdf_path = job_dir / f"{safe_name}.pdf"
            await asyncio.to_thread(to_pdf, notes_md, str(pdf_path))
            jobs[job_id]["files"].append({"name": f"{safe_name}.pdf", "label": "PDF 打印版"})
        except Exception as e:
            # PDF 生成失败，继续处理（GTK 未安装时预期行为）
            pass
        jobs[job_id]["mindmap"] = mindmap_md
        jobs[job_id]["notes_md"] = notes_md  # 返回笔记内容到前端

        # 保存到数据库
        course = get_course_by_name(subject)
        if not course:
            course_id = create_course(subject, color=color)
        else:
            course_id = course['id']

        # 如果有提纲，保存提纲
        outline_id = None
        if outline:
            outline_id = create_outline(course_id, f"{subject}-提纲", outline, {})
            parse_outline_to_sections(outline_id, outline)

        # 保存会话
        save_session(
            session_id=job_id,
            course_id=course_id,
            title=subject,
            outline_id=outline_id,
            materials_count=len(saved_files),
            ai_analysis=raw_text[:1000],  # 保存部分原始文本
            notes_md=notes_md,
            mindmap_md=mindmap_md
        )

        _set(job_id, "done", "✅ 处理完成！")

    except Exception as e:
        # 打印完整堆栈到后端命令行，方便定位
        print("\n" + "=" * 60)
        print(f"[整理流程异常] job_id={job_id}")
        traceback.print_exc()
        print("=" * 60 + "\n")
        _set(job_id, "error", f"❌ 处理出错: {e}")


@app.get("/")
async def index():
    """首页 Landing Page"""
    return FileResponse(str(STATIC_DIR / "home.html"))


@app.get("/library")
async def library():
    """课程笔记库（卡片视图）"""
    return FileResponse(str(STATIC_DIR / "library.html"))


@app.get("/course/{course_id}")
async def course_page(course_id: int):
    """课程主页（标签页布局）"""
    return FileResponse(str(STATIC_DIR / "course.html"))


@app.get("/workspace")
async def workspace():
    """笔记整理工作区"""
    return FileResponse(str(STATIC_DIR / "index.html"))


@app.get("/edit/{session_id}")
async def edit_page(session_id: str):
    """资料编辑页（三栏布局）"""
    return FileResponse(str(STATIC_DIR / "edit.html"))


@app.post("/process")
async def process(
    subject: str = Form(...),
    outline: str = Form(default=""),
    color: str = Form(default="blue"),
    files: list[UploadFile] = File(...),
):
    job_id = uuid.uuid4().hex
    job_dir = UPLOAD_DIR / job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    jobs[job_id] = {"status": "running", "progress": "📤 上传中...", "files": [], "mindmap": "", "notes_md": ""}

    saved: list[Path] = []
    for f in files:
        dest = job_dir / (f.filename or "file")
        async with aiofiles.open(dest, "wb") as out:
            await out.write(await f.read())
        saved.append(dest)

    if IS_SERVERLESS:
        # serverless 平台不保留后台任务，必须在请求内同步处理完
        await _run_pipeline(job_id, subject, outline, saved, color)
    else:
        asyncio.create_task(_run_pipeline(job_id, subject, outline, saved, color))
    return {"job_id": job_id}


@app.get("/status/{job_id}")
async def status(job_id: str):
    # 内存中有就直接返回（本地常驻进程 / 同实例）
    if job_id in jobs:
        return JSONResponse(jobs[job_id])
    # serverless 下轮询请求可能落到新实例，内存查不到 → 回退查数据库。
    # job_id 即 session_id，处理完成会写入 sessions 表。
    session = get_session(job_id)
    if session:
        return JSONResponse({
            "status": "done",
            "progress": "✅ 处理完成！",
            "notes_md": session.get("notes_md", ""),
            "mindmap": session.get("mindmap_md", ""),
            "files": [],
        })
    return JSONResponse({"status": "not_found"})


@app.get("/download/{job_id}/{filename}")
async def download(job_id: str, filename: str):
    path = OUTPUT_DIR / job_id / filename
    if not path.exists():
        return JSONResponse({"error": "file not found"}, status_code=404)
    return FileResponse(path, filename=filename)


@app.post("/api/refine")
async def refine_notes(payload: dict):
    """根据用户指令让 AI 修改笔记，返回修改后的 Markdown。"""
    notes_md = payload.get("notes_md", "")
    instruction = payload.get("instruction", "")
    if not notes_md or not instruction:
        return JSONResponse({"error": "缺少 notes_md 或 instruction"}, status_code=400)
    try:
        from modules.structurer import refine
        new_notes = await asyncio.to_thread(refine, notes_md, instruction)
        return JSONResponse({"notes_md": new_notes})
    except Exception as e:
        print("\n[AI修改异常]")
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/models")
async def list_models():
    """返回可供选择的 AI 模型列表"""
    from modules.structurer import AVAILABLE_MODELS, DEFAULT_MODEL
    return JSONResponse({"models": AVAILABLE_MODELS, "default": DEFAULT_MODEL})


@app.post("/api/chat")
async def chat_api(payload: dict):
    """多模型对话接口。payload: {messages:[{role,content}], model, notes_context}"""
    messages = payload.get("messages", [])
    model = payload.get("model", "")
    notes_context = payload.get("notes_context", "")
    if not messages:
        return JSONResponse({"error": "缺少 messages"}, status_code=400)
    try:
        from modules.structurer import chat
        reply = await asyncio.to_thread(chat, messages, model, notes_context)
        return JSONResponse({"reply": reply})
    except Exception as e:
        print("\n[AI对话异常]")
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/preview/{session_id}/{filename}")
async def preview_file(session_id: str, filename: str):
    """文件预览：PPT/Word 转 PDF（需 LibreOffice），其他类型直接返回原文件。"""
    safe = os.path.basename(filename)
    src = UPLOAD_DIR / session_id / safe
    if not src.exists():
        return JSONResponse({"error": "file not found"}, status_code=404)

    ext = src.suffix.lower()
    if ext in {".ppt", ".pptx", ".doc", ".docx"}:
        # 转换缓存目录
        preview_dir = OUTPUT_DIR / "previews" / session_id
        preview_dir.mkdir(parents=True, exist_ok=True)
        pdf_path = preview_dir / (src.stem + ".pdf")
        if not pdf_path.exists():
            ok = await asyncio.to_thread(_convert_to_pdf, str(src), str(preview_dir))
            if not ok or not pdf_path.exists():
                # 转换失败兜底：提取文字
                return JSONResponse({
                    "fallback": True,
                    "text": _extract_text_fallback(src),
                }, status_code=200)
        return FileResponse(pdf_path, media_type="application/pdf")

    # 其它类型直接给原文件
    return FileResponse(src)


def _convert_to_pdf(src_path: str, out_dir: str) -> bool:
    """用 LibreOffice 将 PPT/Word 转为 PDF。返回是否成功。"""
    import shutil
    import subprocess
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        # 常见 Windows 安装路径兜底
        for cand in (
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ):
            if os.path.exists(cand):
                soffice = cand
                break
    if not soffice:
        print("[预览] 未找到 LibreOffice，无法转换 PDF")
        return False
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, src_path],
            check=True, timeout=120,
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        return True
    except Exception as e:
        print(f"[预览] LibreOffice 转换失败: {e}")
        return False


def _extract_text_fallback(src: Path) -> str:
    """转换失败时提取文字内容用于展示。"""
    ext = src.suffix.lower()
    try:
        if ext in {".ppt", ".pptx"}:
            from pptx import Presentation
            prs = Presentation(str(src))
            pages = []
            for idx, slide in enumerate(prs.slides, 1):
                texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
                pages.append(f"# 第 {idx} 页\n\n" + "\n".join(t for t in texts if t.strip()))
            return "\n\n---\n\n".join(pages)
        elif ext in {".doc", ".docx"}:
            from docx import Document
            doc = Document(str(src))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"无法预览该文件：{e}"
    return "无法预览该文件。"


@app.post("/api/export/docx")
async def export_docx(payload: dict):
    """将前端编辑后的 Markdown 内容导出为 Word 文档。"""
    notes_md = payload.get("notes_md", "")
    subject = payload.get("subject", "课程笔记")
    if not notes_md:
        return JSONResponse({"error": "缺少 notes_md"}, status_code=400)

    safe_name = subject.replace("/", "_").replace("\\", "_") or "课程笔记"
    export_dir = OUTPUT_DIR / "exports"
    export_dir.mkdir(exist_ok=True)
    docx_path = export_dir / f"{safe_name}.docx"

    try:
        from modules.exporter import to_docx
        await asyncio.to_thread(to_docx, notes_md, str(docx_path), subject)
        return FileResponse(
            docx_path,
            filename=f"{safe_name}.docx",
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    except Exception as e:
        print("\n[Word导出异常]")
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/courses/{course_name}/history")
async def get_course_history(course_name: str):
    """获取课程的历史整理记录"""
    course = get_course_by_name(course_name)
    if not course:
        return JSONResponse({"sessions": []})

    sessions = get_course_sessions(course['id'], limit=20)
    return JSONResponse({"sessions": sessions})


@app.get("/api/session/{session_id}/detail")
async def get_session_detail(session_id: str):
    """获取会话详情（用于历史记录加载）"""
    session = get_session(session_id)
    if not session:
        return JSONResponse({"error": "session not found"}, status_code=404)
    # 附带教师提纲文本（供编辑页右栏显示）
    outline_text = ""
    if session.get("outline_id"):
        try:
            outline = get_outline_by_id(session["outline_id"])
            if outline:
                outline_text = outline.get("content", "")
        except Exception:
            pass
    session["outline_text"] = outline_text
    return JSONResponse(session)


# 文件类型辅助
def _file_kind(name: str) -> str:
    s = name.lower()
    if s.rsplit(".", 1)[-1] in {"mp3", "wav", "m4a", "ogg", "flac"}:
        return "audio"
    if s.endswith(".pdf"):
        return "pdf"
    if s.endswith((".ppt", ".pptx")):
        return "ppt"
    if s.endswith((".doc", ".docx")):
        return "word"
    if s.rsplit(".", 1)[-1] in {"jpg", "jpeg", "png", "bmp", "webp", "gif"}:
        return "image"
    if s.endswith((".txt", ".md", ".srt")):
        return "text"
    return "other"


@app.get("/api/session/{session_id}/materials")
async def get_session_materials(session_id: str):
    """列出某次整理上传的原始文件（扫描 uploads/{session_id} 目录）"""
    job_dir = UPLOAD_DIR / session_id
    files = []
    if job_dir.exists():
        for fp in sorted(job_dir.iterdir()):
            if fp.is_file():
                files.append({
                    "name": fp.name,
                    "kind": _file_kind(fp.name),
                    "size": fp.stat().st_size,
                    "url": f"/api/material/{session_id}/{fp.name}",
                })
    return JSONResponse({"materials": files})


@app.get("/api/material/{session_id}/{filename}")
async def serve_material(session_id: str, filename: str):
    """提供原始上传文件用于预览/播放"""
    # 防目录穿越
    safe = os.path.basename(filename)
    path = UPLOAD_DIR / session_id / safe
    if not path.exists():
        return JSONResponse({"error": "file not found"}, status_code=404)
    return FileResponse(path)


@app.put("/api/session/{session_id}/notes")
async def save_session_notes(session_id: str, payload: dict):
    """保存编辑页修改后的笔记内容"""
    notes_md = payload.get("notes_md", "")
    if not get_session(session_id):
        return JSONResponse({"error": "session not found"}, status_code=404)
    update_session_notes(session_id, notes_md)
    return JSONResponse({"ok": True})


@app.patch("/api/session/{session_id}/rename")
async def api_rename_session(session_id: str, payload: dict):
    """重命名整理记录（标题），并同步到笔记正文最上方的一级大标题。"""
    new_title = (payload.get("title") or "").strip()
    if not new_title:
        return JSONResponse({"error": "标题不能为空"}, status_code=400)
    session = get_session(session_id)
    if not session:
        return JSONResponse({"error": "session not found"}, status_code=404)

    # 同步笔记正文里的一级标题（首个 # 标题行）
    notes_md = session.get("notes_md") or ""
    new_notes = _replace_top_heading(notes_md, new_title)
    rename_session(session_id, new_title, new_notes)
    return JSONResponse({"ok": True, "title": new_title})


def _replace_top_heading(notes_md: str, new_title: str) -> str:
    """把 Markdown 笔记中的第一个一级标题(# xxx)替换为新标题；若没有则插入到最前面。"""
    if not notes_md.strip():
        return f"# {new_title}\n"
    lines = notes_md.split("\n")
    for i, line in enumerate(lines):
        if line.lstrip().startswith("# ") and not line.lstrip().startswith("## "):
            lines[i] = f"# {new_title}"
            return "\n".join(lines)
    # 没有一级标题，插到最前
    return f"# {new_title}\n\n" + notes_md


@app.delete("/api/material/{session_id}/{filename}")
async def delete_material(session_id: str, filename: str):
    """删除某次整理上传的原始资料文件。"""
    safe = os.path.basename(filename)
    path = UPLOAD_DIR / session_id / safe
    if not path.exists():
        return JSONResponse({"error": "file not found"}, status_code=404)
    try:
        path.unlink()
        # 同步清理预览缓存
        preview_pdf = OUTPUT_DIR / "previews" / session_id / (Path(safe).stem + ".pdf")
        if preview_pdf.exists():
            preview_pdf.unlink()
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)
    return JSONResponse({"ok": True})


# ============ 课程库管理 API ============

class CourseCreate(BaseModel):
    name: str
    color: str = "blue"


class CourseRename(BaseModel):
    name: str


class CourseColor(BaseModel):
    color: str


@app.get("/api/courses")
async def api_list_courses():
    """列出所有课程及统计数据"""
    return JSONResponse({"courses": list_courses_with_stats()})


@app.post("/api/courses")
async def api_create_course(payload: CourseCreate):
    """新建课程"""
    name = payload.name.strip()
    if not name:
        return JSONResponse({"error": "课程名称不能为空"}, status_code=400)
    existing = get_course_by_name(name)
    if existing:
        return JSONResponse({"error": "课程已存在", "id": existing["id"]}, status_code=409)
    course_id = create_course(name, color=payload.color)
    return JSONResponse({"id": course_id, "name": name, "color": payload.color})


@app.get("/api/courses/id/{course_id}")
async def api_get_course(course_id: int):
    """获取单门课程的详情与会话列表"""
    course = get_course_by_id(course_id)
    if not course:
        return JSONResponse({"error": "course not found"}, status_code=404)
    sessions = get_course_sessions(course_id, limit=100)

    # 扫描每次整理上传目录，汇总该课程下所有原始资料
    materials = []
    for s in sessions:
        job_dir = UPLOAD_DIR / s["id"]
        if job_dir.exists():
            for fp in sorted(job_dir.iterdir()):
                if fp.is_file():
                    materials.append({
                        "name": fp.name,
                        "kind": _file_kind(fp.name),
                        "size": fp.stat().st_size,
                        "session_id": s["id"],
                        "session_title": s["title"],
                        "url": f"/api/material/{s['id']}/{fp.name}",
                    })
    return JSONResponse({"course": course, "sessions": sessions, "materials": materials})


@app.patch("/api/courses/id/{course_id}/rename")
async def api_rename_course(course_id: int, payload: CourseRename):
    """重命名课程"""
    name = payload.name.strip()
    if not name:
        return JSONResponse({"error": "课程名称不能为空"}, status_code=400)
    rename_course(course_id, name)
    return JSONResponse({"ok": True, "name": name})


@app.patch("/api/courses/id/{course_id}/color")
async def api_update_color(course_id: int, payload: CourseColor):
    """更换课程颜色"""
    update_course_color(course_id, payload.color)
    return JSONResponse({"ok": True, "color": payload.color})


@app.delete("/api/courses/id/{course_id}")
async def api_delete_course(course_id: int):
    """删除课程"""
    delete_course(course_id)
    return JSONResponse({"ok": True})


@app.delete("/api/session/{session_id}")
async def api_delete_session(session_id: str):
    """删除单条整理记录"""
    delete_session(session_id)
    return JSONResponse({"ok": True})


@app.post("/api/sessions/batch_delete")
async def api_batch_delete_sessions(payload: dict):
    """批量删除整理记录"""
    ids = payload.get("ids", [])
    for sid in ids:
        delete_session(sid)
    return JSONResponse({"ok": True, "deleted": len(ids)})
